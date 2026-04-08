import os
from pathlib import Path
from typing import List

import pdfplumber

try:
    from docx import Document as DocxDocument
except ImportError:
    DocxDocument = None

try:
    from pptx import Presentation
except ImportError:
    Presentation = None

from .image_preprocessing import ocr_image

SUPPORTED_TEXT_EXTENSIONS = {".txt"}
SUPPORTED_PDF_EXTENSIONS = {".pdf"}
SUPPORTED_IMAGE_EXTENSIONS = {".png", ".jpg", ".jpeg", ".tif", ".tiff", ".bmp"}
SUPPORTED_DOCX_EXTENSIONS = {".docx"}
SUPPORTED_PPTX_EXTENSIONS = {".pptx"}


def read_text_file(path: Path) -> str:
    with path.open("r", encoding="utf-8", errors="ignore") as fh:
        return fh.read().strip()


def read_pdf_file(path: Path) -> str:
    pages: List[str] = []
    with pdfplumber.open(path) as pdf:
        for page in pdf.pages:
            text = page.extract_text() or ""
            pages.append(text)
    return "\n\n".join(pages).strip()


def read_docx_file(path: Path) -> str:
    if DocxDocument is None:
        raise ImportError("请安装 python-docx 以支持 .docx 文件。")
    doc = DocxDocument(path)
    paragraphs = [paragraph.text for paragraph in doc.paragraphs if paragraph.text.strip()]
    return "\n\n".join(paragraphs).strip()


def read_pptx_file(path: Path) -> str:
    if Presentation is None:
        raise ImportError("请安装 python-pptx 以支持 .pptx 文件。")
    prs = Presentation(path)
    slides_text = []
    for slide in prs.slides:
        slide_text = []
        for shape in slide.shapes:
            if hasattr(shape, "text"):
                slide_text.append(shape.text)
        slides_text.append("\n".join(slide_text))
    return "\n\n--- 新幻灯片 ---\n\n".join(slides_text).strip()


def load_document(path: str) -> str:
    path_obj = Path(path)
    if not path_obj.exists():
        raise FileNotFoundError(f"文档不存在: {path}")

    suffix = path_obj.suffix.lower()
    if suffix in SUPPORTED_TEXT_EXTENSIONS:
        return read_text_file(path_obj)
    if suffix in SUPPORTED_PDF_EXTENSIONS:
        return read_pdf_file(path_obj)
    if suffix in SUPPORTED_DOCX_EXTENSIONS:
        return read_docx_file(path_obj)
    if suffix in SUPPORTED_PPTX_EXTENSIONS:
        return read_pptx_file(path_obj)
    if suffix in SUPPORTED_IMAGE_EXTENSIONS:
        return ocr_image(path_obj)

    raise ValueError(f"不支持的文件类型: {suffix}")


def load_image_document(path: str) -> str:
    path_obj = Path(path)
    if path_obj.suffix.lower() not in SUPPORTED_IMAGE_EXTENSIONS:
        raise ValueError(f"文件不是支持的图片类型: {path_obj.suffix}")
    return ocr_image(path_obj)


def split_text(text: str, chunk_size: int = 800, chunk_overlap: int = 100) -> List[str]:
    if not text:
        return []

    paragraphs = [paragraph.strip() for paragraph in text.split("\n\n") if paragraph.strip()]
    chunks: List[str] = []
    current_chunk = []
    current_length = 0

    def flush_chunk():
        nonlocal current_chunk, current_length
        if current_chunk:
            chunks.append("\n\n".join(current_chunk).strip())
            current_chunk = []
            current_length = 0

    for paragraph in paragraphs:
        paragraph_length = len(paragraph)
        if current_length + paragraph_length + 2 <= chunk_size:
            current_chunk.append(paragraph)
            current_length += paragraph_length + 2
        else:
            if current_chunk:
                flush_chunk()
            if paragraph_length <= chunk_size:
                current_chunk.append(paragraph)
                current_length = paragraph_length
            else:
                start = 0
                while start < paragraph_length:
                    end = min(start + chunk_size, paragraph_length)
                    chunks.append(paragraph[start:end].strip())
                    start += chunk_size - chunk_overlap

    flush_chunk()
    return chunks


def load_documents(paths: List[str], chunk_size: int = 800, chunk_overlap: int = 100) -> List[dict]:
    documents: List[dict] = []
    for path in paths:
        text = load_document(path)
        chunks = split_text(text, chunk_size=chunk_size, chunk_overlap=chunk_overlap)
        for idx, chunk in enumerate(chunks):
            documents.append(
                {
                    "content": chunk,
                    "metadata": {
                        "source": os.path.basename(path),
                        "chunk_index": idx,
                        "path": path,
                    },
                }
            )
    return documents

